"""Extract text from uploaded files for pipeline grounding context.

Converts binary file content into plain text that agents can consume.
The extracted text becomes part of grounding_docs, which the executor
injects as brief_context into every agent dispatch payload.

Supported formats: PDF, DOCX, PPTX, TXT, MD, CSV. Unsupported formats
return empty string — the file is still stored in GCS but won't contribute
to grounding context.

Constraint: Libraries (PyPDF2, python-docx, python-pptx) are imported lazily
inside extraction functions. This keeps the import fast for endpoints that
don't need parsing, and fails gracefully if a library is missing.

Constraint: All extraction functions catch Exception broadly and return "".
A corrupt file should not crash a pipeline run — it simply contributes no
text to grounding context. The file's GCS URI is still available for manual
download.

Source reference: SCP briefs (docs/reference/scp-brief-template-extracted.md
in Ops Console) arrive as DOCX or PDF. This parser handles both.
"""

from __future__ import annotations

import io


def extract_text(data: bytes, mime_type: str) -> str:
    """Extract text content from file bytes based on MIME type."""
    if mime_type in ("text/plain", "text/markdown", "text/csv"):
        return data.decode("utf-8", errors="replace")

    if mime_type == "application/pdf":
        return _extract_pdf(data)

    if mime_type in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(data)

    if mime_type in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(data)

    return ""


def _extract_pdf(data: bytes) -> str:
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if texts:
                parts.append(f"Slide {i}:\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except Exception:
        return ""
